"""Text extraction from documents — shared by study_rag ingest and syllabus import.

Supports Markdown/txt/rst plus PDF, DOCX, PPTX. Kept dependency-lazy so importing this
module is cheap; the heavy parsers load only when a matching file is read.
"""

from __future__ import annotations

from pathlib import Path

from polaris_core.logging import get_logger

logger = get_logger(__name__)

TEXT_SUFFIXES = {".md", ".markdown", ".txt", ".rst"}
DOC_SUFFIXES = {".pdf", ".docx", ".pptx"}
SUPPORTED = TEXT_SUFFIXES | DOC_SUFFIXES


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _read_docx(path: Path) -> str:
    import docx2txt

    return docx2txt.process(str(path)) or ""


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n\n".join(p for p in parts if p)


def read_file(path: str | Path) -> str | None:
    """Extract text from a supported file, or None if unreadable/unsupported."""
    path = Path(path)
    suffix = path.suffix.lower()
    try:
        if suffix in TEXT_SUFFIXES:
            return path.read_text(encoding="utf-8")
        if suffix == ".pdf":
            return _read_pdf(path)
        if suffix == ".docx":
            return _read_docx(path)
        if suffix == ".pptx":
            return _read_pptx(path)
    except Exception as exc:  # noqa: BLE001 - skip unreadable files, keep going
        logger.warning("Could not read %s (%s)", path, exc)
    return None


def iter_files(path: str | Path):
    """Yield supported files under a path (recurses directories)."""
    p = Path(path)
    files = [p] if p.is_file() else sorted(p.rglob("*"))
    for f in files:
        if f.is_file() and f.suffix.lower() in SUPPORTED:
            yield f
