"""Ingestion: load notes → split into chunks → embed → store in Chroma.

Supports Markdown/txt/rst plus PDF, DOCX, and PPTX. Ingestion is **incremental**:
each file's content is hashed, so re-ingesting skips unchanged files and re-indexes
(upserts) changed ones instead of piling up duplicate chunks.

Run whenever notes change. Retrieval afterwards is fully offline.
"""

from __future__ import annotations

import hashlib
from pathlib import Path

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from polaris_core.config import Settings, get_settings
from polaris_core.logging import get_logger

from study_rag.store import get_vectorstore

logger = get_logger(__name__)

# Plain-text formats read directly.
TEXT_SUFFIXES = {".md", ".markdown", ".txt", ".rst"}
# Rich formats via loaders (deps installed with the project).
DOC_SUFFIXES = {".pdf", ".docx", ".pptx"}
SUPPORTED = TEXT_SUFFIXES | DOC_SUFFIXES


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _read_docx(path: Path) -> str:
    import docx2txt

    return docx2txt.process(str(path)) or ""


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n\n".join(p for p in parts if p)


def _read_file(path: Path) -> str | None:
    """Extract text from a supported file, or None if unreadable/unsupported."""
    suffix = path.suffix.lower()
    try:
        if suffix in TEXT_SUFFIXES:
            return path.read_text(encoding="utf-8")
        if suffix == ".pdf":
            return _read_pdf(path)
        if suffix == ".docx":
            return _read_docx(path)
        if suffix == ".pptx":
            return _read_pptx(path)
    except Exception as exc:  # noqa: BLE001 - skip unreadable files, keep going
        logger.warning("Skipping %s (%s)", path, exc)
    return None


def _iter_files(path: Path):
    files = [path] if path.is_file() else sorted(path.rglob("*"))
    for f in files:
        if f.is_file() and f.suffix.lower() in SUPPORTED:
            yield f


def _content_hash(text: str) -> str:
    return hashlib.sha1(text.encode("utf-8")).hexdigest()


def ingest_path(target: str | Path, settings: Settings | None = None) -> int:
    """Ingest a file or directory of notes (incrementally). Returns chunks added/updated.

    Unchanged files (matching stored content hash) are skipped; changed files have their
    old chunks removed and replaced.
    """
    settings = settings or get_settings()
    path = Path(target)
    if not path.exists():
        raise FileNotFoundError(f"No such file or directory: {path}")

    store = get_vectorstore(settings)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.rag_chunk_size,
        chunk_overlap=settings.rag_chunk_overlap,
        add_start_index=True,
    )

    total_new = 0
    skipped = 0
    for f in _iter_files(path):
        text = _read_file(f)
        if not text or not text.strip():
            continue
        source = str(f)
        digest = _content_hash(text)

        # Dedupe: compare against any existing chunks for this source.
        existing = store.get(where={"source": source}, include=["metadatas"])
        existing_ids = existing.get("ids", [])
        existing_hashes = {m.get("content_hash") for m in existing.get("metadatas", [])}
        if existing_ids and existing_hashes == {digest}:
            skipped += 1
            continue
        if existing_ids:
            store.delete(ids=existing_ids)  # file changed → drop stale chunks

        doc = Document(
            page_content=text,
            metadata={"source": source, "title": f.stem, "content_hash": digest},
        )
        chunks = splitter.split_documents([doc])
        store.add_documents(chunks)
        total_new += len(chunks)
        logger.info("Indexed %s → %d chunks", f.name, len(chunks))

    if skipped:
        logger.info("Skipped %d unchanged file(s).", skipped)
    if total_new == 0 and skipped == 0:
        logger.warning("No supported files found under %s", path)
    logger.info("Stored %d new/updated chunks in %r", total_new, settings.rag_collection)
    return total_new
